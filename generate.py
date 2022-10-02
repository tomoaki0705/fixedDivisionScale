from cmath import log
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.dml.line import LineFormat
from pptx.shapes.connector import Connector
from pptx import Presentation
from pptx.util import Cm, Pt
import math

defaultLeftPosition = 2
defaultRightPosition = 26
tickerLengthLevel0 = 1.0/2
tickerLengthLevel1 = 0.5/2
tickerLengthLevel2 = 0.2/2

prs = Presentation('A4_template.pptx')
title_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(title_slide_layout)

# title.text = "Hello, World!"
# subtitle.text = "python-pptx was here!"

def computeScale(left, right,scaleOnLine):
    lengthInCentimeter = right - left
    return lengthInCentimeter / scaleOnLine

def drawTicker(begin,inclusiveEnd,scaleGap,slide,y,height=tickerLengthLevel0,left=defaultLeftPosition,right=defaultRightPosition):
    tickerIndex = begin
    scale = computeScale(left,right,math.log10(inclusiveEnd/begin))
    while(tickerIndex <= inclusiveEnd):
        position = scale*(math.log10(tickerIndex/begin))+left
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(position), Cm(y - height), Cm(position), Cm(y + height))
        tickerIndex += scaleGap

def drawTickerInvert(begin,inclusiveEnd,scaleGap,slide,y,height=tickerLengthLevel0,left=defaultLeftPosition,right=defaultRightPosition):
    tickerIndex = begin
    scale = computeScale(left,right,math.log10(inclusiveEnd/begin))
    while(tickerIndex <= inclusiveEnd):
        position = right-scale*(math.log10(tickerIndex/begin))
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(position), Cm(y - height), Cm(position), Cm(y + height))
        tickerIndex += scaleGap

def drawLog10Line(begin,inclusiveEnd,slide,y,height=tickerLengthLevel0,left=defaultLeftPosition,right=defaultRightPosition,indexScale=1):
    scale = computeScale(left, right, math.log10(inclusiveEnd))
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(left), Cm(y), Cm(right), Cm(y))
    for i in range(begin,inclusiveEnd+1):
        scaledIndex = i * indexScale
        # print(f"scaledIndex:{scaledIndex}")
        position=scale*(math.log10(scaledIndex/(begin*indexScale)))+left
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(position), Cm(y-height), Cm(position), Cm(y+height))
        # print(f"position:{position}")
        textBox = slide.shapes.add_textbox(Cm(position-0.5), Cm(y-2), Cm(1), Cm(1))
        paragraph0 = textBox.text_frame
        paragraph0.text = str(scaledIndex)


def patch_connector():
    def get_or_add_ln(self):
        return self._element.spPr.get_or_add_ln()
    Connector.get_or_add_ln = get_or_add_ln
patch_connector()

def drawLog10LineInvert(begin,inclusiveEnd,slide,y,height=tickerLengthLevel0,left=defaultLeftPosition,right=defaultRightPosition,indexScale=1):
    scale = computeScale(left, right, math.log10(inclusiveEnd))
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(left), Cm(y), Cm(right), Cm(y))
    for i in range(begin,inclusiveEnd+1):
        scaledIndex = i * indexScale
        # print(f"scaledIndex:{scaledIndex}")
        position=right - scale*(math.log10(scaledIndex/(begin*indexScale)))
        # line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(position), Cm(y-height), Cm(position), Cm(y+height))
        # line.ln = line.get_or_add_ln
        # lineFormat = LineFormat(line)
        # lineFormat.fill.fore_color.rgb = RGBColor(255, 0, 0)
        # line.color.rgb = RGBColor(255, 0, 0)
        # print(f"position:{position}")
        textBox = slide.shapes.add_textbox(Cm(position-0.5), Cm(y-2), Cm(1), Cm(1))
        paragraph0 = textBox.text_frame
        paragraph0.text = str(scaledIndex)
        # paragraph0.font.color.rgb = RGBColor(255,0,0)
    # for i in range(begin,inclusiveEnd):
    #     scaledIndex = i * indexScale
    #     for j in range(0,10):
    #         fineTics = scaledIndex + 0.1 * j
    #         position=right - scale*(math.log10(fineTics/(begin*indexScale)))
    #         slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(position), Cm(y-0.2), Cm(position), Cm(y+0.2))

def drawLog2Line(begin,end,ticNumber,slide,y,height=tickerLengthLevel0,left=defaultLeftPosition,right=defaultRightPosition,indexScale=1):
    scale = computeScale(left, right, math.log2(end/begin))
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(left), Cm(y), Cm(right), Cm(y))
    # print (f"bgein:{begin}")
    # print (f"end  :{end}")
    leftLog = math.log2(begin)
    # print(f"leftLog    :{leftLog}")
    for i in range(0,ticNumber):
        scaledIndex = i * indexScale + begin
        position=scale*((math.log2(scaledIndex))-leftLog)+left
        # print(f"scaledIndex:{scaledIndex}")
        # print(f"position   :{position}")
        # slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(position), Cm(y-height), Cm(position), Cm(y+height))
        textBox = slide.shapes.add_textbox(Cm(position-0.5), Cm(y-1), Cm(1), Cm(1))
        textFrame0 = textBox.text_frame
        paragraph0 = textFrame0.paragraphs[0]
        paragraph0.text = str(scaledIndex)
        paragraph0.font.size = Pt(8)
        paragraph0.alignment = PP_ALIGN.CENTER
    # for i in range(0,ticNumber-1):
    #     scaledIndex = i * indexScale + begin
    #     for j in range(1,8):
    #         fineTics = (scaledIndex + (j * begin/16)) * 1.0
    #         position=scale*(math.log2(fineTics)-leftLog)+left
    #         # print(f"{fineTics},{position}")
    #         slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(position), Cm(y-0.2), Cm(position), Cm(y+0.2))

# line1=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 1, Cm(2), Cm(1), Cm(2))
# line1=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(1), Cm(1), Cm(3), Cm(3))

verticalOffset = 2.5
offset = verticalOffset
topScale = computeScale(defaultLeftPosition, defaultRightPosition, math.log10(10))
positionOfTwo  = topScale * math.log10(2) + defaultLeftPosition
positionOfFive = topScale * math.log10(5) + defaultLeftPosition
drawLog10Line(1,10,slide,offset)
drawTicker(1,2 ,0.05,slide,offset,height=tickerLengthLevel1,right=positionOfTwo)
drawTicker(1,2 ,0.01,slide,offset,height=tickerLengthLevel2,right=positionOfTwo)
drawTicker(2,5 ,0.1 ,slide,offset,height=tickerLengthLevel1,left=positionOfTwo,right=positionOfFive)
drawTicker(2,5 ,0.02,slide,offset,height=tickerLengthLevel2,left=positionOfTwo,right=positionOfFive)
drawTicker(5,10,0.1 ,slide,offset,height=tickerLengthLevel1,left=positionOfFive)
drawTicker(5,10,0.05,slide,offset,height=tickerLengthLevel2,left=positionOfFive)
offset += verticalOffset
# drawLogLine(1,2,11,2,slide,offset,indexScale=0.1)
# offset += verticalOffset
drawLog2Line(16,160,19,slide,offset,indexScale=8)
drawTicker(16,160,8,slide,offset,height=tickerLengthLevel0)
drawTicker(16,160,4,slide,offset,height=tickerLengthLevel1)
drawTicker(16,160,1,slide,offset,height=tickerLengthLevel2)
offset += verticalOffset
drawLog2Line(160,1600,19,slide,offset,indexScale=80)
drawTicker(160,320 ,16,slide,offset,height=tickerLengthLevel0,right=positionOfTwo)
drawTicker(160,320 ,8 ,slide,offset,height=tickerLengthLevel1,right=positionOfTwo)
drawTicker(160,320 ,2 ,slide,offset,height=tickerLengthLevel2,right=positionOfTwo)
drawTicker(320,800 ,32,slide,offset,height=tickerLengthLevel0,left=positionOfTwo,right=positionOfFive)
drawTicker(320,800 ,16,slide,offset,height=tickerLengthLevel1,left=positionOfTwo,right=positionOfFive)
drawTicker(320,800 ,4 ,slide,offset,height=tickerLengthLevel2,left=positionOfTwo,right=positionOfFive)
drawTicker(800,1600,80,slide,offset,height=tickerLengthLevel0,left=positionOfFive)
drawTicker(800,1600,16,slide,offset,height=tickerLengthLevel1,left=positionOfFive)
drawTicker(800,1600,8 ,slide,offset,height=tickerLengthLevel2,left=positionOfFive)
offset += verticalOffset
# drawLogLine(160,1600,46,2,slide,offset,indexScale=32)
# offset += verticalOffset
drawLog2Line(32,320,19,slide,offset,indexScale=16)
drawTicker(32,160,8,slide,offset,height=tickerLengthLevel0,right=positionOfFive)
drawTicker(32,160,4,slide,offset,height=tickerLengthLevel1,right=positionOfFive)
drawTicker(32,160,1,slide,offset,height=tickerLengthLevel2,right=positionOfFive)
offset += verticalOffset
drawLog2Line(64,640,19,slide,offset,indexScale=32)
offset += verticalOffset
drawLog2Line(128,1280,19,slide,offset,indexScale=64)
offset += verticalOffset

title_slide_layout = prs.slide_layouts[6]
slide2 = prs.slides.add_slide(title_slide_layout)

offset = verticalOffset
drawLog10LineInvert(1,10,slide2,offset)

positionOfTwo  = defaultRightPosition - topScale * math.log10(2)
positionOfFive = defaultRightPosition - topScale * math.log10(5)
drawLog10LineInvert(1,10,slide2,offset)
drawTickerInvert(1,2 ,0.05,slide2,offset,height=tickerLengthLevel1,left=positionOfTwo)
drawTickerInvert(1,2 ,0.01,slide2,offset,height=tickerLengthLevel2,left=positionOfTwo)
drawTickerInvert(2,5 ,0.1 ,slide2,offset,height=tickerLengthLevel1,right=positionOfTwo,left=positionOfFive)
drawTickerInvert(2,5 ,0.02,slide2,offset,height=tickerLengthLevel2,right=positionOfTwo,left=positionOfFive)
drawTickerInvert(5,10,0.1 ,slide2,offset,height=tickerLengthLevel1,right=positionOfFive)
drawTickerInvert(5,10,0.05,slide2,offset,height=tickerLengthLevel2,right=positionOfFive)

divScale=(defaultRightPosition-defaultLeftPosition)
positionOfEight=defaultRightPosition-(divScale * math.log10(8))
offset += verticalOffset
drawLog2Line(2,16,15,slide2,offset,indexScale=1,left=positionOfEight)
drawTicker(2,16,1,slide2,offset,height=tickerLengthLevel0,left=positionOfEight)
offset += verticalOffset
positionOfThree=defaultRightPosition-(divScale * math.log10(32/3))
drawLog2Line(3,32,30,slide2,offset,indexScale=1,left=positionOfThree)
drawTicker(3,32,1,slide2,offset,height=tickerLengthLevel0,left=positionOfThree)
offset += verticalOffset
drawLog2Line(6 ,32,27,slide2,offset,indexScale=1,left=positionOfThree,right=positionOfTwo)
drawLog2Line(32,64,17,slide2,offset,indexScale=2,left=positionOfTwo)
drawTicker(6,32,1,slide2,offset,height=tickerLengthLevel0,left=positionOfThree,right=positionOfTwo)
drawTicker(32,64,2,slide2,offset,height=tickerLengthLevel0,left=positionOfTwo)
drawTicker(32,64,1,slide2,offset,height=tickerLengthLevel1,left=positionOfTwo)
offset += verticalOffset
positionOfFour=defaultRightPosition-(divScale * math.log10(4))
drawLog2Line(12,32 ,21,slide2,offset,indexScale=1,left=positionOfThree,right=positionOfFour)
drawLog2Line(32,64,17,slide2,offset,indexScale=2,left=positionOfFour,right=positionOfTwo)
drawLog2Line(64,128,17,slide2,offset,indexScale=4,left=positionOfTwo)
drawTicker(12,32 ,1,slide2,offset,height=tickerLengthLevel0,left=positionOfThree,right=positionOfFour)
drawTicker(32,64 ,2,slide2,offset,height=tickerLengthLevel0,left=positionOfFour,right=positionOfTwo)
drawTicker(32,64 ,1,slide2,offset,height=tickerLengthLevel1,left=positionOfFour,right=positionOfTwo)
drawTicker(64,128,4,slide2,offset,height=tickerLengthLevel0,left=positionOfTwo)
drawTicker(64,128,1,slide2,offset,height=tickerLengthLevel1,left=positionOfTwo)

prs.save('test.pptx')